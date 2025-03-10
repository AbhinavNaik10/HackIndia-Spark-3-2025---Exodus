from flask import Flask, render_template, request, send_file, jsonify
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from dotenv import load_dotenv
import tempfile
import json
import traceback

# Load environment variables
load_dotenv()

app = Flask(__name__)

# Check for required API keys
def check_api_keys():
    required_keys = {
        'GOOGLE_API_KEY': 'Google API key is required for Gemini AI'
    }
    
    missing_keys = []
    for key, message in required_keys.items():
        if not os.getenv(key) or os.getenv(key) == f'your_{key.lower()}_here':
            missing_keys.append(message)
    
    if missing_keys:
        raise ValueError('\n'.join(missing_keys))

# Configure API keys and models
try:
    check_api_keys()
    genai.configure(api_key=os.getenv('GOOGLE_API_KEY'))
    
    # List available models
    print("Available models:")
    for m in genai.list_models():
        print(f"- {m.name}")
        
except ValueError as e:
    print(f"Error: {str(e)}")
    print("Please update your .env file with valid API keys.")
    exit(1)

def generate_slide_content(topic, num_slides, style):
    """Generate slide content using Gemini AI"""
    try:
        prompt = f"""Create a presentation outline for a {style} presentation about {topic}.
        Generate exactly {num_slides} slides.
        Format the response as a JSON array where each element represents a slide and has this exact structure:
        {{"title": "Slide Title", "content": ["Point 1", "Point 2", "Point 3"]}}
        Make it engaging and professional."""

        # Use gemini-1.0-pro instead of gemini-pro
        model = genai.GenerativeModel('gemini-1.0')
        response = model.generate_content(prompt)
        
        # Add error handling for JSON parsing
        try:
            # Convert response to text and clean it up
            response_text = response.text.strip()
            if not response_text.startswith('['):
                # If response doesn't start with [, try to find the JSON array
                start_idx = response_text.find('[')
                if start_idx != -1:
                    response_text = response_text[start_idx:]
                    end_idx = response_text.rfind(']') + 1
                    if end_idx > 0:
                        response_text = response_text[:end_idx]
            
            content = json.loads(response_text)
            print(f"Generated content: {json.dumps(content, indent=2)}")
            return content
        except json.JSONDecodeError as e:
            print(f"Error parsing Gemini response: {response_text}")
            # Try to format the response as JSON if it's not already
            fallback_content = [
                {"title": topic, "content": ["Generated content could not be parsed as JSON"]},
                {"title": "Error Details", "content": [str(e)]}
            ]
            return fallback_content
            
    except Exception as e:
        print(f"Error in generate_slide_content: {str(e)}")
        print(f"Full traceback: {traceback.format_exc()}")
        raise

def create_presentation(content, style):
    """Create PowerPoint presentation"""
    try:
        prs = Presentation()
        
        # Set slide dimensions
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = content[0]['title']
        subtitle.text = "AI-Generated Presentation"
        
        # Content slides
        for slide_data in content[1:]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.shapes.placeholders[1]
            
            title.text = slide_data['title']
            tf = body.text_frame
            
            for point in slide_data['content']:
                p = tf.add_paragraph()
                p.text = point
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(51, 51, 51)
        
        # Save presentation
        temp_dir = tempfile.gettempdir()
        output_path = os.path.join(temp_dir, 'presentation.pptx')
        prs.save(output_path)
        return output_path
    except Exception as e:
        print(f"Error in create_presentation: {str(e)}")
        print(f"Full traceback: {traceback.format_exc()}")
        raise

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate():
    try:
        topic = request.form['topic']
        num_slides = int(request.form['num_slides'])
        style = request.form['style']
        
        print(f"Generating presentation for topic: {topic}, slides: {num_slides}, style: {style}")
        
        # Generate content
        content = generate_slide_content(topic, num_slides, style)
        print(f"Generated content: {json.dumps(content, indent=2)}")
        
        # Generate presentation
        output_path = create_presentation(content, style)
        print(f"Presentation saved to: {output_path}")
        
        return jsonify({
            'success': True,
            'message': 'Presentation generated successfully',
            'download_url': f'/download/{os.path.basename(output_path)}'
        })
        
    except Exception as e:
        print(f"Error in generate endpoint: {str(e)}")
        print(f"Full traceback: {traceback.format_exc()}")
        return jsonify({
            'success': False,
            'message': str(e)
        }), 500

@app.route('/download/<filename>')
def download(filename):
    try:
        temp_dir = tempfile.gettempdir()
        return send_file(
            os.path.join(temp_dir, filename),
            as_attachment=True,
            download_name=filename
        )
    except Exception as e:
        print(f"Error in download endpoint: {str(e)}")
        print(f"Full traceback: {traceback.format_exc()}")
        return jsonify({
            'success': False,
            'message': str(e)
        }), 500

if __name__ == '__main__':
    app.run(debug=True) 